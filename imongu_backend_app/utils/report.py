from reportlab.lib.pagesizes import A4
from reportlab.pdfgen import canvas
from reportlab.lib import colors
from reportlab.platypus import Table, TableStyle
from reportlab.lib.units import inch
import io
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from io import BytesIO

def generate_pdf(user, role_name, goals, okrs):
    buffer = BytesIO()
    p = canvas.Canvas(buffer, pagesize=A4)
    p.setFillColorRGB(0.97, 0.98, 0.98)  
    p.rect(0, 0, A4[0], A4[1], fill=1)

    p.setFont("Helvetica-Bold", 10)
    p.setFillColorRGB(1, 0.627, 0)  
    p.drawString(50, 800, "User: ")
    p.setFillColorRGB(0, 0, 0)  
    p.drawString(80, 800, user.username)  

    # Draw Email label and value
    p.setFont("Helvetica-Bold", 10)
    p.setFillColorRGB(1, 0.627, 0) 
    p.drawString(50, 780, "Email:  ")
    p.setFillColorRGB(0, 0, 0)  
    p.drawString(80, 780, user.email) 

    p.setFont("Helvetica-Bold", 10)
    p.setFillColorRGB(1, 0.627, 0)  
    p.drawString(50, 760, "Role: ")
    p.setFillColorRGB(0, 0, 0)  
    p.drawString(80, 760, role_name)
    p.setLineWidth(1)
    p.setStrokeColorRGB(0.95, 0.49, 0.05)  
    p.line(50, 750, 550, 750)

    p.setFillColorRGB(1, 0.63, 0.26)  
    p.setFont("Helvetica-Bold", 16)
    p.drawString(250, 730, "Goals")
    p.setFont("Helvetica", 12)

    goals_data = [['Session', 'Title', 'Avg Gain', 'Created At']]
    for goal in goals:
        goals_data.append([goal.session, goal.title, goal.average_gain, goal.created_at.strftime('%Y-%m-%d')])

    goals_table = Table(goals_data, colWidths=[1.5 * inch, 2 * inch, 1.5 * inch, 2 * inch])
    goals_table.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (-1, 0), colors.transparent),  
        ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
        ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
        ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
        ('FONTSIZE', (0, 0), (-1, 0), 12),
        ('BOTTOMPADDING', (0, 0), (-1, 0), 10),
        ('BACKGROUND', (0, 1), (-1, -1), colors.white),  
        ('GRID', (0, 0), (-1, -1), 0.5, colors.black),
        ('BACKGROUND', (0, 0), (-1, 0), colors.Color(1, 0.63, 0.26)),  
    ]))

    goals_y_position = 700
    goals_table.wrapOn(p, 50, goals_y_position)
    goals_table.drawOn(p, 50, goals_y_position - len(goals_data) * 18)

    okrs_y_position = goals_y_position - len(goals_data) * 20 - 60
    p.setFont("Helvetica-Bold", 16)
    p.setFillColorRGB(1, 0.63, 0.26)  
    p.drawString(250, okrs_y_position, "OKRs")
    p.setFont("Helvetica", 12)

    okrs_data = [['Session', 'Title', 'Avg Gain', 'Goal Title']]
    for okr_item in okrs:
        okrs_data.append([okr_item.session, okr_item.title, okr_item.average_gain, okr_item.goal_id.title])

    okrs_table = Table(okrs_data, colWidths=[1.5 * inch, 2 * inch, 1.5 * inch, 2 * inch])
    okrs_table.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (-1, 0), colors.transparent),
        ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
        ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
        ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
        ('FONTSIZE', (0, 0), (-1, 0), 12),
        ('BOTTOMPADDING', (0, 0), (-1, 0), 10),
        ('BACKGROUND', (0, 1), (-1, -1), colors.white),
        ('GRID', (0, 0), (-1, -1), 0.5, colors.black),
        ('BACKGROUND', (0, 0), (-1, 0), colors.Color(1, 0.63, 0.26)),  
    ]))

    okrs_table.wrapOn(p, 50, okrs_y_position - 40)
    okrs_table.drawOn(p, 50, okrs_y_position - len(okrs_data) * 18 - 20)

    # Finalize PDF
    p.showPage()
    p.save()
    pdf_content = buffer.getvalue()
    buffer.close()

    return pdf_content, "user_report.pdf"


def generate_ppt(user, role_name, goals, okrs):
    prs = Presentation()
    slide_layout = prs.slide_layouts[1]  
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "User Report"
    text_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(2))
    text_frame = text_box.text_frame
    text_frame.word_wrap = True

    user_info = [
        ("Username: ", user.username),
        ("Email: ", user.email),
        ("Role: ", role_name)
    ]
    
    for label, value in user_info:
        paragraph = text_frame.add_paragraph()
    
        run_label = paragraph.add_run()
        run_label.text = label
        run_label.font.size = Pt(18)
        run_label.font.bold = True
        run_label.font.color.rgb = RGBColor(255, 160, 67)  

        run_value = paragraph.add_run()
        run_value.text = value
        run_value.font.size = Pt(18)
        run_value.font.color.rgb = RGBColor(0, 0, 0)  
        paragraph.alignment = PP_ALIGN.LEFT

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title_shape = slide.shapes.title
    title_shape.text = "Goals"
    title_text_frame = title_shape.text_frame
    title_paragraph = title_text_frame.paragraphs[0]
    title_paragraph.font.size = Pt(24)
    title_paragraph.font.bold = True
    title_paragraph.font.color.rgb = RGBColor(255, 160, 67)
    
    goals_data = [['Session', 'Title', 'Avg Gain', 'Created At']]
    for goal in goals:
        goals_data.append([
            goal.session, 
            goal.title, 
            str(goal.average_gain), 
            goal.created_at.strftime('%Y-%m-%d')
        ])
    
    rows = len(goals_data)
    cols = len(goals_data[0])
    table = slide.shapes.add_table(rows, cols, Inches(1), Inches(1.5), Inches(8), Inches(3)).table

    for i in range(cols):
        table.columns[i].width = Inches(2)
    
    for row_idx, row_data in enumerate(goals_data):
        for col_idx, cell_value in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(cell_value)
            cell.text_frame.paragraphs[0].font.size = Pt(12)
            
            # Setting alignment for each cell
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            
            if row_idx == 0:  # Header row
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(255, 160, 41)
                cell.text_frame.paragraphs[0].font.bold = True
                cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
                cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title_shape = slide.shapes.title
    title_shape.text = "OKRs"
    title_text_frame = title_shape.text_frame
    title_paragraph = title_text_frame.paragraphs[0]
    title_paragraph.font.size = Pt(24)
    title_paragraph.font.bold = True
    title_paragraph.font.color.rgb = RGBColor(255, 160, 67)
    
    okrs_data = [['Session', 'Title', 'Avg Gain', 'Goal Title']]
    for okr_item in okrs:
        okrs_data.append([
            okr_item.session, 
            okr_item.title, 
            str(okr_item.average_gain), 
            okr_item.goal_id.title
        ])
    
    rows = len(okrs_data)
    cols = len(okrs_data[0])
    table = slide.shapes.add_table(rows, cols, Inches(1), Inches(1.5), Inches(8), Inches(3)).table

    for i in range(cols):
        table.columns[i].width = Inches(2)
    
    for row_idx, row_data in enumerate(okrs_data):
        for col_idx, cell_value in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(cell_value)
            cell.text_frame.paragraphs[0].font.size = Pt(12)
            
            # Setting alignment for each cell
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

            if row_idx == 0:  # Header row
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(255, 160, 41)
                cell.text_frame.paragraphs[0].font.bold = True
                cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
                cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)

    ppt_buffer = BytesIO()
    prs.save(ppt_buffer)
    ppt_buffer.seek(0)

    return ppt_buffer.getvalue(), "user_report.pptx"

